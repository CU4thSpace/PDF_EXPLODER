import os
import sys
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pathlib import Path
from pdf2image import convert_from_path
import tempfile
from tkinter import filedialog, Tk, messagebox
from PIL import Image

def convert_pdf_to_images(pdf_path, output_folder, dpi=200):
    """Converts PDF pages to images and returns a list of saved image paths."""
    poppler_path = "/opt/homebrew/bin"  # <-- set this explicitly
    # Convert once to get every page
    images = convert_from_path(pdf_path, dpi=dpi, poppler_path=poppler_path)
    total_pages = len(images)           # e.g. 132
    num_digits  = len(str(total_pages)) # → 3  (because '132' has 3 digits) to fix sorting error
    image_paths = []
    for page_num, img in enumerate(images, start=1):
        filename = f"pdf_page_{page_num:0{num_digits}d}.png"
        img_path = os.path.join(output_folder, filename)
        img.save(img_path, "PNG")
        image_paths.append(img_path)
    return image_paths

def create_ppt_from_images(folder_path, output_file="output.pptx"):
    prs = Presentation()
    # Set to 16:9 aspect ratio
    prs.slide_width = Inches(20)
    prs.slide_height = Inches(11.25)
    blank_slide_layout = prs.slide_layouts[6]  # blank layout
    
    image_extensions = ('.png', '.jpg', '.jpeg', '.bmp', '.gif')
    image_files = sorted([
        f for f in os.listdir(folder_path)
        if f.lower().endswith(image_extensions)
    ])
    
    if not image_files:
        print("No image files found in the folder.")
        return
    
    for image_file in image_files:
        image_path = os.path.join(folder_path, image_file)
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Set slide background to black
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black
        
        # Get image dimensions to calculate proper sizing
        with Image.open(image_path) as img:
            img_width, img_height = img.size
        
        # Calculate aspect ratios
        slide_aspect = float(prs.slide_width) / float(prs.slide_height)  # 16:9 ≈ 1.78
        img_aspect = img_width / img_height
        
        # Calculate dimensions to fit image within slide while maintaining aspect ratio
        if img_aspect > slide_aspect:
            # Image is wider than slide - fit to width, letterbox top/bottom
            new_width = prs.slide_width
            new_height = prs.slide_width / img_aspect
            left = Inches(0)
            top = (prs.slide_height - new_height) / 2
        else:
            # Image is taller than slide - fit to height, pillarbox left/right
            new_height = prs.slide_height
            new_width = prs.slide_height * img_aspect
            left = (prs.slide_width - new_width) / 2
            top = Inches(0)
        
        # Add the image with calculated dimensions and position
        pic = slide.shapes.add_picture(image_path, left, top, new_width, new_height)
    
    prs.save(output_file)
    print(f"Presentation saved as: {output_file}")

# === USAGE ===
if __name__ == "__main__":
    # If arguments are provided (drag-and-drop):
    if len(sys.argv) > 1:
        input_path = sys.argv[1]
        desktop_path = Path.home() / "Desktop"
        output_name = sys.argv[2] if len(sys.argv) > 2 else "output.pptx"
        output = str(desktop_path / output_name)
    else:
        root = Tk()
        root.withdraw()  # Hide the root window
        
        # Ask user: PDF file or folder of images?
        choice = messagebox.askyesno(
            title="Choose Input Type",
            message="Select YES to choose a PDF file, or NO to select a folder of images."
        )
        
        if choice:
            # Select PDF file
            input_path = filedialog.askopenfilename(
                title="Select a PDF file",
                filetypes=[("PDF files", "*.pdf")]
            )
        else:
            # Select folder of images
            input_path = filedialog.askdirectory(title="Select a folder of images")
        
        if not input_path:
            print("No file or folder selected. Exiting.")
            exit()
        
        output_name = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            initialfile="output.pptx",
            title="Save PowerPoint As",
            filetypes=[("PowerPoint files", "*.pptx")]
        )
        
        if not output_name:
            print("No output filename selected. Exiting.")
            exit()
        
        output = output_name  # Use selected output path
    
    if not os.path.exists(input_path):
        print("Error: File or folder does not exist.")
        exit()
    
    if input_path.lower().endswith(".pdf"):
        # Convert PDF to images in a temp directory
        with tempfile.TemporaryDirectory() as temp_dir:
            print("Converting PDF to images...")
            convert_pdf_to_images(input_path, temp_dir)
            create_ppt_from_images(temp_dir, output)
    elif os.path.isdir(input_path):
        create_ppt_from_images(input_path, output)
    else:
        print("Error: Please enter a valid folder path or PDF file.")
